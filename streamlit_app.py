
import streamlit as st
from qdrant_client import QdrantClient
from qdrant_client.http.models import VectorParams, Distance, PointStruct
from sentence_transformers import SentenceTransformer
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from langchain_openai import ChatOpenAI
import openai
import os

st.set_page_config(page_title="Housing Disrepair QA", layout="wide")

# Load API keys with fallback to environment
openai.api_key = st.secrets["openai"].get("api_key") or os.environ.get("OPENAI_API_KEY")
qdrant_url = st.secrets["qdrant"]["url"]
qdrant_key = st.secrets["qdrant"]["api_key"]

# Connect to Qdrant
qdrant_client = QdrantClient(url=qdrant_url, api_key=qdrant_key)
embedding_model = SentenceTransformer("all-MiniLM-L6-v2")

# Ensure collection exists
COLLECTION_NAME = "paterson_docs"
if not qdrant_client.collection_exists(collection_name=COLLECTION_NAME):
    qdrant_client.create_collection(
        collection_name=COLLECTION_NAME,
        vectors_config=VectorParams(size=384, distance=Distance.COSINE)
    )

# Functions
def embed_text_chunks(chunks):
    return [
        PointStruct(id=i, vector=embedding_model.encode(c).tolist(), payload={"text": c})
        for i, c in enumerate(chunks)
    ]

def upsert_documents(text_chunks):
    points = embed_text_chunks(text_chunks)

def query_qdrant(query_text, top_k=5):
    hits = []
    try:
        vector = embedding_model.encode(query_text).tolist()
    try:
        hits = qdrant_client.query_points(
    except Exception as e:
        st.error(f"Qdrant query failed: {e}")
        hits = []
            collection_name=COLLECTION_NAME,
            vector=vector,
            limit=top_k
        )
    except Exception as e:
        st.warning(f"Qdrant search failed: {e}")
    return [hit.payload["text"] for hit in hits if "text" in hit.payload]

def extract_text_from_file(file):
    if file.name.endswith(".pdf"):
        reader = PdfReader(file)
        return " ".join([page.extract_text() or "" for page in reader.pages])
    elif file.name.endswith(".docx"):
        doc = docx.Document(file)
        return "\n".join([p.text for p in doc.paragraphs])
    elif file.name.endswith(".pptx"):
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    elif file.name.endswith(".txt"):
        return file.read().decode("utf-8")
    else:
        return ""

# UI
st.title("📂 Upload & Ask – Housing Disrepair Assistant")

uploaded_file = st.file_uploader("Upload a document (.pdf, .docx, .pptx, .txt)", type=["pdf", "docx", "pptx", "txt"])
if uploaded_file:
    text = extract_text_from_file(uploaded_file)
    if text.strip():
        st.success("✅ Document uploaded and indexed into Qdrant.")
        chunks = [text[i:i+500] for i in range(0, len(text), 500)]
        upsert_documents(chunks)
    else:
        st.warning("⚠️ No readable text found in the uploaded file.")

st.divider()
st.subheader("💬 Ask a Question")

user_query = st.chat_input("Ask me about housing disrepair...")
context_docs = []
context_text = ''
answer = None
context_docs = []
answer = None
if user_query:
    with st.chat_message("user"):
        st.markdown(user_query)

    context_docs = query_qdrant(user_query) or []
    context_text = "\n".join(context_docs) if context_docs else ""
    context_text = "\n".join(context_docs)

    if context_docs:
        prompt = f"Answer the question based on the following documents:\n\n{context_text}\n\nQuestion: {user_query}"
    else:
        prompt = user_query

    llm = ChatOpenAI(temperature=0, model="gpt-3.5-turbo", api_key=st.secrets["openai"].get("api_key") or os.environ.get("OPENAI_API_KEY"))
try:
    if context_text:
        prompt = f"Answer the question based on the following documents:\n\n{context_text}\n\nQuestion: {user_query}"
    else:
        prompt = user_query
    prompt = prompt[:4000]
    try:
        answer = llm.invoke(prompt)
    except Exception as e:
        st.error(f"OpenAI call failed: {e}")
        answer = None
    if context_docs:
        prompt = f"Answer the question based on the following documents:\n\n{context_text}\n\nQuestion: {user_query}"
    else:
        prompt = user_query
    prompt = prompt[:4000]  # truncate prompt if too long
    answer = llm.invoke(prompt)
    prompt = prompt[:4000]  # truncate prompt if too long
    answer = llm.invoke(prompt)
    prompt = prompt[:4000]  # truncate prompt if too long
    answer = llm.invoke(prompt)
except Exception as e:
    st.error(f"OpenAI request failed: {str(e)}")
    answer = None
    st.error(f"OpenAI request failed: {str(e)}")
    answer = None
    st.error(f"OpenAI request failed: {str(e)}")
    st.error(f"OpenAI request failed: {str(e)}")
    answer = None
    st.error(f"OpenAI request failed: {str(e)}")
    answer = None

    with st.chat_message("assistant"):
        if context_docs:
            st.markdown("**Context (from Qdrant):**")
            for i, doc in enumerate(context_docs, 1):
                st.markdown(f"**{i}.** {doc}")
            st.markdown("---")
    if answer is not None:
        st.markdown(f"**Answer:** {answer.content}")
